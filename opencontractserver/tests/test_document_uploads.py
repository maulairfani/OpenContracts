import io

from django.conf import settings
from django.contrib.auth import get_user_model
from django.test import TestCase
from docx import Document as DocxDocument
from graphene.test import Client
from graphql_relay import from_global_id, to_global_id
from openpyxl import Workbook
from pptx import Presentation

from config.graphql.schema import schema
from opencontractserver.corpuses.models import Corpus
from opencontractserver.documents.models import Document as DocumentModel
from opencontractserver.types.enums import PermissionTypes
from opencontractserver.utils.files import base_64_encode_bytes
from opencontractserver.utils.permissioning import set_permissions_for_obj_to_user

User = get_user_model()


class TestContext:
    def __init__(self, user):
        self.user = user


class UploadDocumentMutationTestCase(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(
            username="testuser", password="testpassword"
        )
        self.client = Client(schema, context_value=TestContext(self.user))
        self.mutation = """
            mutation UploadDocument(
                $file: String!,
                $filename: String!,
                $title: String!,
                $description: String!,
                $customMeta: GenericScalar!,
                $addToCorpusId: ID,
                $makePublic: Boolean!
            ) {
                uploadDocument(
                    base64FileString: $file,
                    filename: $filename,
                    title: $title,
                    description: $description,
                    customMeta: $customMeta,
                    addToCorpusId: $addToCorpusId,
                    makePublic: $makePublic
                ) {
                    ok
                    message
                    document {
                        id
                        title
                        fileType
                    }
                }
            }
        """

    def generate_file_content(self, file_type):
        if file_type == "pdf":
            return b"%PDF-1.0\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj 2 0 obj<</Type/Pages/Count 0>>endobj\nxref\n0 3\n0000000000 65535 f \n0000000010 00000 n \n0000000059 00000 n \ntrailer<</Size 3/Root 1 0 R>>startxref\n104\n%%EOF"  # noqa: E501
        elif file_type == "docx":
            buffer = io.BytesIO()
            doc = DocxDocument()
            doc.add_paragraph("This is a test DOCX file.")
            doc.save(buffer)
            return buffer.getvalue()
        elif file_type == "xlsx":
            buffer = io.BytesIO()
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "This is a test XLSX file."
            wb.save(buffer)
            return buffer.getvalue()
        elif file_type == "pptx":
            buffer = io.BytesIO()
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = "This is a test PPTX file."
            prs.save(buffer)
            return buffer.getvalue()
        elif file_type == "txt":
            return b"This is a text file."
        else:
            raise ValueError(f"Unsupported file type for generation: {file_type}")

    def test_upload_document(self):
        allowed_mimetypes = set(settings.ALLOWED_DOCUMENT_MIMETYPES)

        file_types_to_test = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "txt": "text/plain",
        }

        for file_ext, mime_type in file_types_to_test.items():
            with self.subTest(file_type=file_ext):
                file_content = self.generate_file_content(file_ext)
                base64_content = base_64_encode_bytes(file_content)

                result = self.client.execute(
                    self.mutation,
                    variables={
                        "file": base64_content,
                        "filename": f"test.{file_ext}",
                        "title": f"Test {file_ext.upper()}",
                        "description": f"A test {file_ext.upper()} file",
                        "makePublic": True,
                        "customMeta": {},
                    },
                )

                self.assertIsNone(
                    result.get("errors"), f"GraphQL errors found for {file_ext}"
                )

                upload_result = result["data"]["uploadDocument"]
                is_allowed = mime_type in allowed_mimetypes

                if is_allowed:
                    self.assertTrue(
                        upload_result["ok"],
                        f"Upload should succeed for allowed type {mime_type}",
                    )
                    self.assertEqual(
                        upload_result["message"],
                        "Success",
                        f"Success message expected for {mime_type}",
                    )
                    self.assertIsNotNone(
                        upload_result["document"],
                        f"Document data should exist for successful upload of {mime_type}",
                    )
                    self.assertEqual(
                        upload_result["document"]["title"],
                        f"Test {file_ext.upper()}",
                        f"Incorrect title for {mime_type}",
                    )
                    self.assertEqual(
                        upload_result["document"]["fileType"],
                        mime_type,
                        f"Incorrect fileType for {mime_type}",
                    )

                    doc_id_str = upload_result["document"]["id"]
                    doc_pk = from_global_id(doc_id_str)[1]
                    doc = DocumentModel.objects.get(id=doc_pk)
                    self.assertEqual(doc.title, f"Test {file_ext.upper()}")
                    self.assertEqual(doc.file_type, mime_type)
                    self.assertEqual(doc.creator, self.user)
                    self.assertTrue(doc.is_public)

                    if file_ext == "txt":
                        self.assertTrue(
                            bool(doc.txt_extract_file),
                            f"txt_extract_file should exist for {file_ext}",
                        )
                    else:
                        self.assertTrue(
                            bool(doc.pdf_file), f"pdf_file should exist for {file_ext}"
                        )
                else:
                    self.assertFalse(
                        upload_result["ok"],
                        f"Upload should fail for disallowed type {mime_type}",
                    )
                    self.assertIn(
                        "Unallowed filetype",
                        upload_result["message"],
                        f"Expected 'Unallowed filetype' message for {mime_type}",
                    )
                    self.assertIsNone(
                        upload_result["document"],
                        f"Document data should be null for failed upload of {mime_type}",
                    )

    def test_upload_document_to_corpus_requires_edit_permission(self):
        """
        Test that uploading a document to a corpus requires EDIT permission.
        Users without EDIT permission should be denied.
        """
        # Create another user who owns a corpus
        corpus_owner = User.objects.create_user(
            username="corpus_owner", password="testpassword"
        )
        corpus = Corpus.objects.create(
            title="Owner's Corpus",
            creator=corpus_owner,
        )
        set_permissions_for_obj_to_user(corpus_owner, corpus, [PermissionTypes.ALL])

        # Generate a simple PDF file
        file_content = self.generate_file_content("pdf")
        base64_content = base_64_encode_bytes(file_content)
        corpus_global_id = to_global_id("CorpusType", corpus.id)

        # Test 1: User with NO permissions on corpus should be denied
        result = self.client.execute(
            self.mutation,
            variables={
                "file": base64_content,
                "filename": "test.pdf",
                "title": "Unauthorized Upload",
                "description": "This should fail",
                "makePublic": False,
                "customMeta": {},
                "addToCorpusId": corpus_global_id,
            },
        )

        self.assertIsNone(result.get("errors"), "GraphQL errors found")
        upload_result = result["data"]["uploadDocument"]
        self.assertFalse(upload_result["ok"], "Upload should fail without permission")
        self.assertIn(
            "permission",
            upload_result["message"].lower(),
            "Error message should mention permission",
        )
        self.assertIsNone(
            upload_result["document"],
            "Document should not be created without permission",
        )

        # Test 2: User with only READ permission should still be denied
        set_permissions_for_obj_to_user(self.user, corpus, [PermissionTypes.READ])

        result = self.client.execute(
            self.mutation,
            variables={
                "file": base64_content,
                "filename": "test.pdf",
                "title": "Read-Only Upload",
                "description": "This should also fail",
                "makePublic": False,
                "customMeta": {},
                "addToCorpusId": corpus_global_id,
            },
        )

        self.assertIsNone(result.get("errors"), "GraphQL errors found")
        upload_result = result["data"]["uploadDocument"]
        self.assertFalse(
            upload_result["ok"], "Upload should fail with only READ permission"
        )
        self.assertIn(
            "permission",
            upload_result["message"].lower(),
            "Error message should mention permission",
        )

        # Test 3: User with UPDATE permission should succeed
        # Note: EDIT is an alias for UPDATE in permission checks, but set_permissions
        # uses UPDATE to grant the actual permission
        set_permissions_for_obj_to_user(self.user, corpus, [PermissionTypes.UPDATE])

        result = self.client.execute(
            self.mutation,
            variables={
                "file": base64_content,
                "filename": "test.pdf",
                "title": "Authorized Upload",
                "description": "This should succeed",
                "makePublic": False,
                "customMeta": {},
                "addToCorpusId": corpus_global_id,
            },
        )

        self.assertIsNone(result.get("errors"), "GraphQL errors found")
        upload_result = result["data"]["uploadDocument"]
        self.assertTrue(
            upload_result["ok"], "Upload should succeed with UPDATE permission"
        )
        self.assertIsNotNone(
            upload_result["document"],
            "Document should be created with UPDATE permission",
        )

        # Test 4: Corpus owner should always be able to upload
        owner_client = Client(schema, context_value=TestContext(corpus_owner))

        result = owner_client.execute(
            self.mutation,
            variables={
                "file": base64_content,
                "filename": "test.pdf",
                "title": "Owner Upload",
                "description": "Owner should always succeed",
                "makePublic": False,
                "customMeta": {},
                "addToCorpusId": corpus_global_id,
            },
        )

        self.assertIsNone(result.get("errors"), "GraphQL errors found")
        upload_result = result["data"]["uploadDocument"]
        self.assertTrue(upload_result["ok"], "Owner should be able to upload")
        self.assertIsNotNone(
            upload_result["document"],
            "Document should be created for owner",
        )

    def tearDown(self):
        # TestCase wraps each test in a transaction that gets rolled back,
        # so explicit cleanup is optional. We do clean documents for the
        # original test which runs outside our permission test.
        from opencontractserver.documents.models import DocumentPath

        # DocumentPath has PROTECT on documents, so delete paths first
        DocumentPath.objects.all().delete()
        # Documents are now linked via DocumentPath only - M2M not needed
        DocumentModel.objects.all().delete()
        Corpus.objects.all().delete()
